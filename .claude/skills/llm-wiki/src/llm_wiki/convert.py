"""Doc -> Markdown conversion, pure Python (no markitdown / pandoc / poppler).

Dispatches on file extension. Each converter aims for clean, readable
Markdown that the wiki can treat as ordinary text -- not pixel-perfect
fidelity. Tables are emitted as GitHub-flavored Markdown.
"""
import os
import re


def slugify(name: str) -> str:
    """lowercase, hyphens, no spaces -- matches SCHEMA.md file-name rule."""
    base = os.path.splitext(os.path.basename(name))[0]
    base = base.lower()
    base = re.sub(r"[^a-z0-9]+", "-", base)
    return base.strip("-") or "source"


def sha256_body(text: str) -> str:
    import hashlib

    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _pdf_to_md(path: str) -> str:
    import fitz  # PyMuPDF

    parts = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc):
            txt = page.get_text("text").strip()
            if txt:
                parts.append("<!-- page {} -->\n\n{}".format(i + 1, txt))
    return "\n\n".join(parts)


def _pptx_to_md(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides):
        out.append("## Slide {}".format(i + 1))
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = ("".join(r.text for r in para.runs) or para.text).strip()
                    if text:
                        out.append(text)
            if getattr(shape, "has_table", False) and shape.has_table:
                out.append(_table_to_md(
                    [[c.text for c in row.cells] for row in shape.table.rows]
                ))
        out.append("")
    return "\n\n".join(out).strip()


def _docx_to_md(path: str) -> str:
    import docx

    doc = docx.Document(path)
    out = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()
        if style.startswith("heading"):
            level = "".join(ch for ch in style if ch.isdigit()) or "2"
            out.append("{} {}".format("#" * min(int(level) + 1, 6), text))
        else:
            out.append(text)
    for table in doc.tables:
        out.append(_table_to_md([[c.text for c in row.cells] for row in table.rows]))
    return "\n\n".join(out).strip()


def _xlsx_to_md(path: str) -> str:
    import pandas as pd

    sheets = pd.read_excel(path, sheet_name=None)
    out = []
    for name, df in sheets.items():
        out.append("## {}".format(name))
        out.append(df.fillna("").to_markdown(index=False))
    return "\n\n".join(out).strip()


def _csv_to_md(path: str) -> str:
    import pandas as pd

    df = pd.read_csv(path)
    return df.fillna("").to_markdown(index=False)


def _html_to_md(path: str) -> str:
    from markdownify import markdownify

    with open(path, "r", encoding="utf-8", errors="replace") as fh:
        return markdownify(fh.read()).strip()


def _table_to_md(rows) -> str:
    rows = [[(c or "").replace("\n", " ").strip() for c in r] for r in rows if r]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join(["---"] * width) + " |"
    body = ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return "\n".join([header, sep] + body)


_CONVERTERS = {
    ".pdf": _pdf_to_md,
    ".pptx": _pptx_to_md,
    ".docx": _docx_to_md,
    ".xlsx": _xlsx_to_md,
    ".xls": _xlsx_to_md,
    ".csv": _csv_to_md,
    ".html": _html_to_md,
    ".htm": _html_to_md,
}

SUPPORTED = sorted(_CONVERTERS)


def convert(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    fn = _CONVERTERS.get(ext)
    if fn is None:
        raise ValueError(
            "unsupported extension {!r}; supported: {}".format(ext, ", ".join(SUPPORTED))
        )
    return fn(path)
